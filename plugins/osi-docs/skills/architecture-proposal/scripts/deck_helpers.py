# -*- coding: utf-8 -*-
# deck_helpers.py — 提案pptx用の共通ヘルパー（python-pptx）
# import して使う:  from deck_helpers import *
# 16:9 (13.333 x 7.5 in)。日本語フォントは Noto Sans CJK JP（QA時のLibreOffice描画でCJK可）。
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- ブランド配色（CAIO資料準拠 / EDIT可） ----
DARK="111827"; RED="B91C1C"; CARD="1F2937"; CARD2="3D3D4F"
LIGHT="F7F7F7"; LRED="FEF2F2"; TMED="4B5563"; TLITE="9CA3AF"; WHITE="FFFFFF"
GREEN="27AE60"; ORANGE="E67E22"
# 担当色（構築=赤 / 相手社AI設計=青 / 発注元=グレー）
AIOSI=RED; PARTNER="1D4ED8"; CLIENT="374151"; SHARED="7C3AED"
FONT="Noto Sans CJK JP"

def C(h): return RGBColor.from_string(h)
def tint(h,f=0.9):
    r=int(h[0:2],16);g=int(h[2:4],16);b=int(h[4:6],16)
    return "%02X%02X%02X"%(int(r+(255-r)*f),int(g+(255-g)*f),int(b+(255-b)*f))

def new_deck():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    return prs

def slide(prs,bg=WHITE):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=C(bg); r.line.fill.background(); r.shadow.inherit=False
    return s

def rect(s,x,y,w,h,fill,line=None,lw=0.75,rounded=False):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                           Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=C(fill)
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=C(line); shp.line.width=Pt(lw)
    shp.shadow.inherit=False; return shp

def txt(s,x,y,w,h,text,size,color,bold=False,align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,font=FONT,italic=False,wrap=True):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=wrap
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0; tf.vertical_anchor=anchor
    for i,ln in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
        r.font.name=font; r.font.color.rgb=C(color)
    return tb

def head(s,kicker,title,color=RED,SW=13.333):
    rect(s,0,0,SW,0.09,color)
    txt(s,0.7,0.34,SW-1.4,0.3,kicker,12.5,color,bold=True)
    txt(s,0.7,0.62,SW-1.0,0.7,title,26,DARK,bold=True)

def pagenum(s,n): txt(s,12.45,7.05,0.7,0.3,str(n),10,TLITE,align=PP_ALIGN.RIGHT)

def chip(s,x,y,w,h,fill,text,tcolor=WHITE,size=9.5,bold=True):
    rect(s,x,y,w,h,fill,rounded=True)
    txt(s,x,y,w,h,text,size,tcolor,bold=bold,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

def add_image(s,path,x=0.92,y=1.62,width=11.5):
    s.shapes.add_picture(path,Inches(x),Inches(y),width=Inches(width))

# QA: PDF化→画像化（はみ出し/重なり/略語未定義/課題対応漏れをサブエージェントで点検）
#   soffice --headless --convert-to pdf deck.pptx
#   pdftoppm -jpeg -r 140 deck.pdf slide
